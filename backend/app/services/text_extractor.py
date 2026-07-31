import os
import re
import io
import logging
from typing import Dict, Any, List
from PIL import Image

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("TextExtractor")

try:
    import fitz  # PyMuPDF high-precision text extractor
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

import docx
from pptx import Presentation

# Optional OCR engines
HAS_EASYOCR = False
_easyocr_reader = None

import importlib.util
HAS_EASYOCR = importlib.util.find_spec("easyocr") is not None

HAS_GEMINI_CLIENT = False
try:
    from google import genai
    HAS_GEMINI_CLIENT = True
except ImportError:
    HAS_GEMINI_CLIENT = False

from app.core.config import settings


WATERMARK_PATTERNS = [
    re.compile(r'page\s+\d+\s+of\s+\d+', re.IGNORECASE),
    re.compile(r'page\s+no\.?\s*[/:\d]*', re.IGNORECASE),
    re.compile(r'notes\s+by\s+multi\s+atoms', re.IGNORECASE),
    re.compile(r'subscribe\s+.*youtube.*', re.IGNORECASE),
    re.compile(r'multi\s+atoms\s+youtube', re.IGNORECASE),
    re.compile(r'^\s*page\s+\d+\s*$', re.IGNORECASE),
    re.compile(r'^\s*multi\s+atoms\s*$', re.IGNORECASE),
    re.compile(r'youtube\s+channel\s+for\s+more\s+subjects', re.IGNORECASE)
]


def clean_watermark_lines(text: str) -> str:
    if not text:
        return ""
    lines = text.split('\n')
    cleaned = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        is_watermark = False
        for pat in WATERMARK_PATTERNS:
            if pat.search(stripped):
                is_watermark = True
                break
        if not is_watermark:
            cleaned.append(line)
    return "\n".join(cleaned)


def _get_easyocr_reader():
    global _easyocr_reader
    if _easyocr_reader is None and HAS_EASYOCR:
        try:
            import easyocr
            _easyocr_reader = easyocr.Reader(['en'], gpu=False)
        except Exception as e:
            logger.warning(f"EasyOCR initialization failed: {e}")
            _easyocr_reader = None
    return _easyocr_reader


def ocr_fitz_page(page) -> str:
    # 1. Try Gemini Vision OCR if API key is configured with model fallback
    gemini_key = getattr(settings, "GEMINI_API_KEY", None) or os.environ.get("GEMINI_API_KEY")
    if HAS_GEMINI_CLIENT and gemini_key:
        models_to_try = ["gemini-flash-latest", "gemini-3.5-flash", "gemini-3.1-flash-lite", "gemini-3.6-flash"]
        try:
            client = genai.Client(api_key=gemini_key)
            pix = page.get_pixmap(dpi=150)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            for m in models_to_try:
                try:
                    res = client.models.generate_content(
                        model=m,
                        contents=["Extract all study notes text from this page. Omit watermarks and channel subscribe headers/footers.", img]
                    )
                    if res and res.text:
                        cleaned = clean_watermark_lines(res.text.strip())
                        if cleaned:
                            return cleaned
                except Exception as e_m:
                    logger.warning(f"Gemini Vision model '{m}' failed: {e_m}")
                    continue
        except Exception as e:
            logger.warning(f"Gemini Vision OCR fallback exception: {e}")

    # 2. Try EasyOCR offline engine fallback
    reader = _get_easyocr_reader()
    if reader is not None:
        try:
            pix = page.get_pixmap(dpi=150)
            import numpy as np
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            results = reader.readtext(np.array(img), detail=0)
            raw_ocr = "\n".join(results)
            return clean_watermark_lines(raw_ocr)
        except Exception as e:
            logger.warning(f"EasyOCR page processing exception: {e}")

    return ""


class TextExtractorService:
    @staticmethod
    def _clean_repetitive_watermarks(page_texts: List[str]) -> List[str]:
        cleaned_pages = []
        for p in page_texts:
            cleaned_pages.append(clean_watermark_lines(p))
        return cleaned_pages

    def extract_text(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Extracts raw text and metadata from PDF, DOCX, PPTX, or TXT files.
        Uses PyMuPDF (fitz) for layout-aware PDF text extraction with OCR fallback for scanned pages.
        """
        ext = file_type.upper()
        raw_text = ""
        pages_count = 1

        try:
            if ext == "PDF":
                page_texts = []
                if HAS_FITZ:
                    try:
                        doc = fitz.open(file_path)
                        pages_count = len(doc)
                        logger.info(f"[TextExtractor] Processing PDF with {pages_count} pages...")
                        for idx, page in enumerate(doc):
                            blocks = page.get_text("blocks")
                            block_texts = []
                            for b in blocks:
                                b_text = b[4].strip()
                                if b_text:
                                    cleaned_b = clean_watermark_lines(b_text)
                                    if cleaned_b:
                                        block_texts.append(cleaned_b)
                            
                            p_str = "\n".join(block_texts).strip()

                            # If digital text extraction yields < 40 chars, run OCR for scanned pages
                            if len(p_str) < 40:
                                ocr_text = ocr_fitz_page(page)
                                if ocr_text:
                                    p_str = ocr_text

                            if p_str:
                                page_texts.append(f"[Page {idx+1}]\n" + p_str)
                    except Exception as e:
                        logger.error(f"[TextExtractor] fitz extraction failed: {e}")
                        page_texts = []

                if not page_texts and HAS_PYPDF:
                    reader = PdfReader(file_path)
                    pages_count = len(reader.pages)
                    for idx, page in enumerate(reader.pages):
                        t = page.extract_text() or ""
                        cleaned_t = clean_watermark_lines(t)
                        if cleaned_t.strip():
                            page_texts.append(f"[Page {idx+1}]\n" + cleaned_t)

                cleaned_texts = self._clean_repetitive_watermarks(page_texts)
                raw_text = "\n\n".join(cleaned_texts)

            elif ext == "DOCX":
                doc = docx.Document(file_path)
                paras = [clean_watermark_lines(p.text) for p in doc.paragraphs if p.text.strip()]
                raw_text = "\n".join([p for p in paras if p.strip()])
                pages_count = max(1, len(paras) // 5)

            elif ext == "PPTX":
                prs = Presentation(file_path)
                pages_count = len(prs.slides)
                slide_texts = []
                for idx, slide in enumerate(prs.slides):
                    st = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            cleaned = clean_watermark_lines(shape.text)
                            if cleaned:
                                st.append(cleaned)
                    if st:
                        slide_texts.append(f"[Slide {idx+1}]\n" + "\n".join(st))
                raw_text = "\n\n".join(slide_texts)

            else:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    raw_text = f.read()
                raw_text = clean_watermark_lines(raw_text)
                pages_count = max(1, len(raw_text) // 2000)

        except Exception as e:
            raw_text = f"Text extraction warning: {str(e)}"
            pages_count = 1

        chunk_count = max(1, len(raw_text) // 1500) if raw_text else 0

        return {
            "text": raw_text,
            "pages": pages_count,
            "chunks_count": chunk_count
        }

text_extractor = TextExtractorService()


